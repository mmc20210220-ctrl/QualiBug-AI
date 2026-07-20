"""
QualiBug 昆山落地介绍 PPT 完美版生成脚本
优化方向：文案升级、视觉层次、数据冲击力、设计精致度
"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

# === 配色方案 ===
BLUE = RGBColor(0x0F, 0x4C, 0x81)       # 主色-深蓝
TEAL = RGBColor(0x00, 0xA6, 0xA6)       # 辅色-青
ORANGE = RGBColor(0xFF, 0x8A, 0x00)     # 强调-橙
DARK = RGBColor(0x1F, 0x29, 0x37)       # 正文深色
GRAY = RGBColor(0x6B, 0x72, 0x80)       # 副标题灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # 白色
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)   # 浅背景
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)   # 卡片背景
BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xF8) # 蓝色浅底
TEAL_LIGHT = RGBColor(0xE6, 0xF7, 0xF7) # 青色浅底
ORANGE_LIGHT = RGBColor(0xFF, 0xF3, 0xE0) # 橙色浅底

FONT = 'Noto Sans CJK SC'

# === 尺寸常量 (EMU) ===
SLIDE_W = 12192000
SLIDE_H = 6858000
MARGIN_L = Emu(548640)    # ~0.6 inch
MARGIN_T = Emu(320040)
CONTENT_W = Emu(11094720) # 可用宽度

# === 工具函数 ===
def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, radius=Emu(91440)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # 设置圆角
    shape.adjustments[0] = 0.04
    return shape

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=Pt(12), bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name=FONT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=Pt(12), color=DARK, spacing=Pt(8), bullet_char='•'):
    """添加带项目符号的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'{bullet_char} {item}'
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
    return txBox

def add_footer(slide, page_num):
    """添加统一页脚"""
    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(6446520), Emu(11094720), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    line.line.fill.background()
    # 左侧文字
    add_text_box(slide, MARGIN_L, Emu(6519672), Emu(7315200), Emu(164592),
                 'QualiBug AI智能测试平台｜昆山｜2026年7月', Pt(9), False, GRAY)
    # 右侧页码
    add_text_box(slide, Emu(11155680), Emu(6519672), Emu(502920), Emu(164592),
                 f'{page_num:02d}', Pt(10), False, GRAY, PP_ALIGN.RIGHT)

def add_slide_title(slide, title, subtitle=None):
    """添加统一的页面标题区"""
    add_text_box(slide, MARGIN_L, MARGIN_T, Emu(9144000), Emu(457200),
                 title, Pt(26), True, BLUE)
    if subtitle:
        add_text_box(slide, Emu(576072), Emu(786384), Emu(9144000), Emu(274320),
                     subtitle, Pt(11), False, GRAY)

def add_flow_arrow(slide, left, top, width=Emu(365760)):
    """添加流程箭头"""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, Emu(274320))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    shape.line.fill.background()
    return shape

def add_flow_step(slide, left, top, text, color, width=Emu(1463040), height=Emu(548640)):
    """添加流程步骤块"""
    shape = add_rounded_rect(slide, left, top, width, height, color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return shape


# === 从原始PPT提取图片 ===
orig_prs = Presentation(r'c:\Users\Test\Downloads\QualiBug_昆山落地材料包\QualiBug_昆山落地介绍_202607.pptx')
images = {}
for i, slide in enumerate(orig_prs.slides, 1):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            try:
                img = shape.image
                images[f'slide{i}'] = {
                    'blob': img.blob,
                    'content_type': img.content_type,
                    'ext': 'jpg' if 'jpeg' in img.content_type else 'png'
                }
            except:
                pass

# === 创建新PPT ===
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank layout

# ============================================================
# SLIDE 1 - 封面
# ============================================================
slide = prs.slides.add_slide(blank_layout)

# 左侧背景色块
bg = add_rect(slide, 0, 0, Emu(7543800), SLIDE_H, WHITE)

# 右侧图片
if 'slide1' in images:
    img_stream = io.BytesIO(images['slide1']['blob'])
    slide.shapes.add_picture(img_stream, Emu(7543800), 0, Emu(4648200), SLIDE_H)

# 左侧装饰竖线
accent = add_rect(slide, Emu(731520), Emu(1097280), Pt(4), Emu(2743200), BLUE)

# 主标题
add_text_box(slide, Emu(822960), Emu(1097280), Emu(5486400), Emu(548640),
             'QualiBug', Pt(40), True, BLUE)

# 副标题
add_text_box(slide, Emu(822960), Emu(1691640), Emu(5486400), Emu(411480),
             'AI智能测试平台', Pt(24), True, DARK)

# 核心描述
add_text_box(slide, Emu(822960), Emu(2286000), Emu(5943600), Emu(731520),
             '让AI深度理解企业业务，自主执行测试、发现真实缺陷，\n输出完整可审计的证据链——从"测不了"到"测得透"。',
             Pt(14), False, DARK)

# 标签栏
tag_bar = add_rounded_rect(slide, Emu(822960), Emu(3337560), Emu(4114800), Emu(365760), BLUE)
tf = tag_bar.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.text = '人工智能 × 企业软件测试 × 工业数字化'
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER

# 公司信息
add_text_box(slide, Emu(822960), Emu(5029200), Emu(5486400), Emu(548640),
             '拟设公司：昆山观迹知因智能科技有限公司', Pt(10.5), False, GRAY)
add_text_box(slide, Emu(822960), Emu(5303520), Emu(5486400), Emu(548640),
             '沟通对象：政务中心 / OPC社区 / 招商与试点企业', Pt(10.5), False, GRAY)

# ============================================================
# SLIDE 2 - 为什么需要QualiBug
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, '为什么需要QualiBug',
                '企业软件复杂度指数级增长，传统测试已触及天花板——超过80%的缺陷在测试阶段无法被发现')

# 三列卡片
card_w = Emu(3383280)
card_h = Emu(4297680)
card_y = Emu(1234440)
gap = Emu(365760)
x1 = Emu(548640)
x2 = x1 + card_w + gap
x3 = x2 + card_w + gap

# 卡片1 - 企业侧痛点
add_rounded_rect(slide, x1, card_y, card_w, card_h, BLUE_LIGHT, RGBColor(0xBF, 0xDB, 0xFE))
add_text_box(slide, x1 + Emu(182880), card_y + Emu(137160), Emu(2926080), Emu(274320),
             '企业侧痛点', Pt(15), True, BLUE)
add_bullet_list(slide, x1 + Emu(182880), card_y + Emu(502920), Emu(2926080), Emu(3429000), [
    '需求、接口、UI、数据库等资料高度分散',
    '测试人员理解成本高，重复劳动占比超60%',
    '中小企业缺少专业测试团队，质量靠"碰运气"',
    '版本迭代快，回归测试覆盖严重不足',
], Pt(11), DARK, Pt(10))

# 卡片2 - 传统自动化痛点
add_rounded_rect(slide, x2, card_y, card_w, card_h, TEAL_LIGHT, RGBColor(0x99, 0xF6, 0xE4))
add_text_box(slide, x2 + Emu(182880), card_y + Emu(137160), Emu(2926080), Emu(274320),
             '传统自动化痛点', Pt(15), True, TEAL)
add_bullet_list(slide, x2 + Emu(182880), card_y + Emu(502920), Emu(2926080), Emu(3429000), [
    '依赖人工预设脚本，编写周期长',
    '只能回归已知路径，无法主动发现未知缺陷',
    '每次需求变化都需大量维护脚本',
    '缺乏业务语义理解，误报率高',
], Pt(11), DARK, Pt(10))

# 卡片3 - AI测试现状痛点
add_rounded_rect(slide, x3, card_y, card_w, card_h, ORANGE_LIGHT, RGBColor(0xFF, 0xE0, 0xB2))
add_text_box(slide, x3 + Emu(182880), card_y + Emu(137160), Emu(2926080), Emu(274320),
             'AI测试现状痛点', Pt(15), True, ORANGE)
add_bullet_list(slide, x3 + Emu(182880), card_y + Emu(502920), Emu(2926080), Emu(3429000), [
    '多数工具停留在"生成用例文本"层面',
    '缺少真实环境执行能力和证据链',
    '无法形成缺陷修复后的回归闭环',
    '不能对接企业真实系统和数据',
], Pt(11), DARK, Pt(10))

add_footer(slide, 2)

# ============================================================
# SLIDE 3 - 产品闭环
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, '产品闭环',
                '从企业资料导入到缺陷证据，再到自动回归验证——全链路自动化')

# 核心差异说明
diff_box = add_rounded_rect(slide, Emu(548640), Emu(1143000), Emu(11094720), Emu(457200), RGBColor(0xEF, 0xF6, 0xFF))
tf = diff_box.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '核心差异：不是生成测试用例文本，而是打通 "理解 → 执行 → 发现 → 证据 → 回归" 全链路闭环'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = BLUE
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER

# 6步流程
steps = ['资料导入', '业务建模', '场景生成', '自动执行', '缺陷证据', '回归验证']
step_colors = [BLUE, BLUE, TEAL, TEAL, ORANGE, ORANGE]
step_w = Emu(1554480)
step_h = Emu(502920)
step_y = Emu(1920240)
arrow_w = Emu(274320)
start_x = Emu(548640)

for i, (step_text, color) in enumerate(zip(steps, step_colors)):
    x = start_x + i * (step_w + arrow_w)
    add_flow_step(slide, x, step_y, step_text, color, step_w, step_h)
    if i < len(steps) - 1:
        add_flow_arrow(slide, x + step_w + Emu(27432), step_y + Emu(114300), Emu(219456))

# 三列信息框
info_y = Emu(2926080)
info_w = Emu(3474720)
info_h = Emu(2834640)
info_gap = Emu(320040)
ix1 = Emu(548640)
ix2 = ix1 + info_w + info_gap
ix3 = ix2 + info_w + info_gap

# 输入资料
add_rounded_rect(slide, ix1, info_y, info_w, info_h, CARD_BG)
add_text_box(slide, ix1 + Emu(182880), info_y + Emu(137160), Emu(3017520), Emu(274320),
             '📥 输入资料', Pt(14), True, BLUE)
add_bullet_list(slide, ix1 + Emu(182880), info_y + Emu(502920), Emu(3017520), Emu(2194560), [
    'PRD / 需求文档',
    '接口文档 / 数据库结构',
    'UI原型 / 设计图',
    '历史Bug / 测试报告',
    '测试账号与环境配置',
], Pt(11), DARK, Pt(8))

# 输出结果
add_rounded_rect(slide, ix2, info_y, info_w, info_h, CARD_BG)
add_text_box(slide, ix2 + Emu(182880), info_y + Emu(137160), Emu(3017520), Emu(274320),
             '📤 输出结果', Pt(14), True, TEAL)
add_bullet_list(slide, ix2 + Emu(182880), info_y + Emu(502920), Emu(3017520), Emu(2194560), [
    '真实缺陷与风险点定位',
    '完整复现步骤',
    '请求响应 / 截图 / 日志',
    '数据库状态变化记录',
    '修复后自动回归验证结果',
], Pt(11), DARK, Pt(8))

# 最终价值
add_rounded_rect(slide, ix3, info_y, info_w, info_h, CARD_BG)
add_text_box(slide, ix3 + Emu(182880), info_y + Emu(137160), Emu(3017520), Emu(274320),
             '🎯 最终价值', Pt(14), True, ORANGE)
add_bullet_list(slide, ix3 + Emu(182880), info_y + Emu(502920), Emu(3017520), Emu(2194560), [
    '测试成本降低 60%+',
    '缺陷发现率显著提升',
    '形成可审计完整证据链',
    '交付周期大幅缩短',
    '沉淀行业测试知识资产',
], Pt(11), DARK, Pt(8))

add_footer(slide, 3)

# ============================================================
# SLIDE 4 - 核心能力
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, '核心能力',
                '以"业务行为场景"为主线，覆盖接口、UI、数据和性能四大维度')

# 右侧图片
if 'slide4' in images:
    img_stream = io.BytesIO(images['slide4']['blob'])
    slide.shapes.add_picture(img_stream, Emu(7543800), Emu(1143000), Emu(4251960), Emu(4297680))

# 左上 - 业务理解与场景建模
cap1_y = Emu(1143000)
add_rounded_rect(slide, Emu(548640), cap1_y, Emu(3383280), Emu(1920240), BLUE_LIGHT, RGBColor(0xBF, 0xDB, 0xFE))
add_text_box(slide, Emu(731520), cap1_y + Emu(137160), Emu(3017520), Emu(274320),
             '业务理解与场景建模', Pt(14), True, BLUE)
add_bullet_list(slide, Emu(731520), cap1_y + Emu(457200), Emu(3017520), Emu(1371600), [
    '识别业务流程、规则、角色、权限',
    '覆盖正常路径、异常路径和边界条件',
    '将业务行为映射到接口、页面和数据状态',
], Pt(11), DARK, Pt(8))

# 右上 - 自动执行与缺陷识别
add_rounded_rect(slide, Emu(4114800), cap1_y, Emu(3108960), Emu(1920240), TEAL_LIGHT, RGBColor(0x99, 0xF6, 0xE4))
add_text_box(slide, Emu(4297680), cap1_y + Emu(137160), Emu(2743200), Emu(274320),
             '自动执行与缺陷识别', Pt(14), True, TEAL)
add_bullet_list(slide, Emu(4297680), cap1_y + Emu(457200), Emu(2743200), Emu(1371600), [
    '接口调用、页面操作、数据库校验',
    '检测业务规则、契约和数据一致性问题',
    '输出可复现、可验证、可回归的缺陷证据',
], Pt(11), DARK, Pt(8))

# 下方 - 长期积累
cap2_y = Emu(3383280)
add_rounded_rect(slide, Emu(548640), cap2_y, Emu(6675120), Emu(1920240), ORANGE_LIGHT, RGBColor(0xFF, 0xE0, 0xB2))
add_text_box(slide, Emu(731520), cap2_y + Emu(137160), Emu(6217920), Emu(274320),
             '长期积累与持续演进', Pt(14), True, ORANGE)
add_bullet_list(slide, Emu(731520), cap2_y + Emu(457200), Emu(6217920), Emu(1371600), [
    '企业系统行为模型——越用越懂你的业务',
    '行业测试场景库——跨项目复用测试经验',
    '缺陷知识与证据模板——标准化质量输出',
], Pt(11), DARK, Pt(8))
# 补充说明
add_text_box(slide, Emu(731520), cap2_y + Emu(1554480), Emu(6217920), Emu(274320),
             '从"一次性测试工具"升级为"持续演进的质量基础设施"', Pt(11), True, ORANGE)

add_footer(slide, 4)

# ============================================================
# SLIDE 5 - 当前进展与知识产权规划
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, '当前进展与知识产权规划',
                '核心能力已验证，知识产权布局清晰')

# 左侧 - 当前进展
prog_y = Emu(1143000)
prog_w = Emu(5303520)
prog_h = Emu(3200400)
add_rounded_rect(slide, Emu(548640), prog_y, prog_w, prog_h, BLUE_LIGHT, RGBColor(0xBF, 0xDB, 0xFE))
add_text_box(slide, Emu(731520), prog_y + Emu(137160), Emu(4846320), Emu(274320),
             '✅ 当前进展', Pt(15), True, BLUE)
add_bullet_list(slide, Emu(731520), prog_y + Emu(502920), Emu(4846320), Emu(1828800), [
    '已完成平台核心模块开发（资料解析、场景生成、任务执行、证据采集）',
    '前端管理、任务执行和结果展示形成完整闭环',
    '基准测试已发现 9 个真实缺陷（131-Bug 标准基准）',
    '支持多源异构资料自动解析与业务建模',
], Pt(11), DARK, Pt(8))

# 重点指标
metric_box = add_rounded_rect(slide, Emu(731520), prog_y + Emu(2377440), Emu(4846320), Emu(548640), WHITE)
tf = metric_box.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '🎯 当前重点：持续提升真实缺陷发现率，目标覆盖更多行业场景'
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = BLUE
p.font.name = FONT

# 右侧 - 知识产权规划
ip_x = Emu(6217920)
add_rounded_rect(slide, ip_x, prog_y, prog_w, prog_h, TEAL_LIGHT, RGBColor(0x99, 0xF6, 0xE4))
add_text_box(slide, ip_x + Emu(182880), prog_y + Emu(137160), Emu(4846320), Emu(274320),
             '📋 知识产权规划', Pt(15), True, TEAL)
add_bullet_list(slide, ip_x + Emu(182880), prog_y + Emu(502920), Emu(4846320), Emu(2560320), [
    '公司成立后申请QualiBug平台软件著作权',
    '围绕多源资料解析、业务场景建模、智能测试执行、缺陷证据链申请发明专利',
    '申请公司名称、产品名称及核心标识商标',
    '建立代码、算法、文档和商业秘密管理机制',
], Pt(11), DARK, Pt(8))

# 底部强调条
bottom_bar = add_rounded_rect(slide, Emu(1097280), Emu(4800600), Emu(10012680), Emu(640080), RGBColor(0xEF, 0xF6, 0xFF))
tf = bottom_bar.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '技术壁垒明确：多源资料语义解析 + 业务行为建模 + 自主执行引擎 + 证据链闭环，四项核心能力均具备专利保护价值'
p.font.size = Pt(11)
p.font.bold = False
p.font.color.rgb = BLUE
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER

add_footer(slide, 5)

# ============================================================
# SLIDE 6 - 昆山落地计划
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, '昆山落地计划',
                '优先依托OPC社区和数字经济生态，快速完成试点验证与标杆打造')

# 左侧图片
if 'slide6' in images:
    img_stream = io.BytesIO(images['slide6']['blob'])
    slide.shapes.add_picture(img_stream, Emu(548640), Emu(1143000), Emu(4023360), Emu(3977640))

# 右上 - 政务中心OPC专窗
r_x = Emu(4937760)
r_w = Emu(3200400)
r_h = Emu(1920240)
add_rounded_rect(slide, r_x, Emu(1143000), r_w, r_h, BLUE_LIGHT, RGBColor(0xBF, 0xDB, 0xFE))
add_text_box(slide, r_x + Emu(182880), Emu(1280160), Emu(2834640), Emu(274320),
             '第一站：政务中心OPC专窗', Pt(13), True, BLUE)
add_bullet_list(slide, r_x + Emu(182880), Emu(1600200), Emu(2834640), Emu(1371600), [
    '咨询注册、社区入驻、工位和注册地址',
    '核实创业、人才、租房和社保补贴条件',
    '对接社区运营负责人，建立长期合作',
], Pt(10.5), DARK, Pt(6))

# 右上2 - 优先社区
r2_x = Emu(8458200)
add_rounded_rect(slide, r2_x, Emu(1143000), r_w, r_h, TEAL_LIGHT, RGBColor(0x99, 0xF6, 0xE4))
add_text_box(slide, r2_x + Emu(182880), Emu(1280160), Emu(2834640), Emu(274320),
             '优先社区：全球数字创新港', Pt(13), True, TEAL)
add_bullet_list(slide, r2_x + Emu(182880), Emu(1600200), Emu(2834640), Emu(1371600), [
    '贴近软件/SaaS和AI应用方向',
    '便于对接数字经济企业与应用场景',
    '适合形成产品试点和标杆案例',
], Pt(10.5), DARK, Pt(6))

# 右下 - 试点路径
path_y = Emu(3383280)
path_w = Emu(6720840)
add_rounded_rect(slide, r_x, path_y, path_w, Emu(1920240), ORANGE_LIGHT, RGBColor(0xFF, 0xE0, 0xB2))
add_text_box(slide, r_x + Emu(182880), path_y + Emu(137160), Emu(6263640), Emu(274320),
             '🚀 试点路径（6个月）', Pt(13), True, ORANGE)
add_bullet_list(slide, r_x + Emu(182880), path_y + Emu(502920), Emu(6263640), Emu(1371600), [
    '注册公司并入驻社区（第1个月）',
    '完成1-3家真实企业系统验证（第2-4个月）',
    '形成真实Bug证据和回归报告（第4-5个月）',
    '沉淀行业测试场景模板，打造可复制标杆案例（第5-6个月）',
], Pt(10.5), DARK, Pt(6))

add_footer(slide, 6)

# ============================================================
# SLIDE 7 - 本次沟通希望获得的支持
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_slide_title(slide, '本次沟通希望获得的支持',
                '不是只要工位，更重要的是试点场景和产业资源')

# 三列卡片
card_w = Emu(3474720)
card_h = Emu(3429000)
card_y = Emu(1143000)
gap = Emu(320040)
cx1 = Emu(548640)
cx2 = cx1 + card_w + gap
cx3 = cx2 + card_w + gap

# 基础支持
add_rounded_rect(slide, cx1, card_y, card_w, card_h, BLUE_LIGHT, RGBColor(0xBF, 0xDB, 0xFE))
add_text_box(slide, cx1 + Emu(182880), card_y + Emu(137160), Emu(3017520), Emu(274320),
             '🏢 基础支持', Pt(14), True, BLUE)
add_bullet_list(slide, cx1 + Emu(182880), card_y + Emu(502920), Emu(3017520), Emu(2743200), [
    'OPC社区入驻和创业工位',
    '工商注册地址和注册指导',
    '代理记账、银行开户、税务登记等服务确认',
], Pt(11), DARK, Pt(10))

# 技术与政策支持
add_rounded_rect(slide, cx2, card_y, card_w, card_h, TEAL_LIGHT, RGBColor(0x99, 0xF6, 0xE4))
add_text_box(slide, cx2 + Emu(182880), card_y + Emu(137160), Emu(3017520), Emu(274320),
             '⚙️ 技术与政策支持', Pt(14), True, TEAL)
add_bullet_list(slide, cx2 + Emu(182880), card_y + Emu(502920), Emu(3017520), Emu(2743200), [
    'AI算力、模型和技术资源对接',
    '软件、AI、科技创业政策指导',
    '软著、专利、商标和科技项目申报辅导',
], Pt(11), DARK, Pt(10))

# 最关键支持
add_rounded_rect(slide, cx3, card_y, card_w, card_h, ORANGE_LIGHT, RGBColor(0xFF, 0xE0, 0xB2))
add_text_box(slide, cx3 + Emu(182880), card_y + Emu(137160), Emu(3017520), Emu(274320),
             '🌟 最关键支持', Pt(14), True, ORANGE)
add_bullet_list(slide, cx3 + Emu(182880), card_y + Emu(502920), Emu(3017520), Emu(2743200), [
    '对接昆山本地企业（制造、软件、SaaS等）试点',
    '提供真实企业系统验证机会',
    '帮助形成首批标杆案例和客户背书',
], Pt(11), DARK, Pt(10))

# 底部CTA
cta_box = add_rounded_rect(slide, Emu(1097280), Emu(5029200), Emu(10012680), Emu(548640), BLUE)
tf = cta_box.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '🎯 一句话请求：请帮助 QualiBug 对接 OPC 全球数字创新港，引入 1-3 家昆山本地企业（制造/软件/SaaS均可）作为首批试点——让我们在昆山做出看得见的标杆案例。'
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER

add_footer(slide, 7)

# === 保存 ===
output_path = r'c:\Users\Test\Downloads\QualiBug_昆山落地材料包\QualiBug_昆山落地介绍_202607_v2.pptx'
prs.save(output_path)
print(f'完美版PPT已生成: {output_path}')
print(f'共 {len(prs.slides)} 页')
