import sys
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

prs = Presentation(r'c:\Users\Test\Downloads\QualiBug_昆山落地材料包\QualiBug_昆山落地介绍_202607.pptx')

out = open(r'd:\QualiBug-AI\QualiBug-AI-main\_ppt_content_utf8.txt', 'w', encoding='utf-8')

def p(s=''):
    out.write(s + '\n')

p(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
p(f'Total slides: {len(prs.slides)}')
p('='*80)

for i, slide in enumerate(prs.slides, 1):
    p(f'\n--- Slide {i} ---')
    p(f'Layout: {slide.slide_layout.name}')
    for shape in slide.shapes:
        p(f'  Shape: {shape.shape_type}, name={shape.name}, pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    font_info = ''
                    if para.runs:
                        r = para.runs[0]
                        try:
                            color = r.font.color.rgb if r.font.color and r.font.color.type else None
                        except:
                            color = None
                        font_info = f' [font={r.font.name}, size={r.font.size}, bold={r.font.bold}, color={color}]'
                    p(f'    P: "{text}"{font_info}')
        if shape.has_table:
            table = shape.table
            p(f'    Table: {len(table.rows)}x{len(table.columns)}')
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                p(f'      Row {row_idx}: {cells}')
        if hasattr(shape, 'image'):
            try:
                img = shape.image
                p(f'    Image: {img.content_type}, size={len(img.blob)} bytes')
            except:
                pass

out.close()
print('Done - output written to _ppt_content_utf8.txt')
