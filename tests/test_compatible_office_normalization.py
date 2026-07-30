from __future__ import annotations

import hashlib
import io

import pytest

from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.builtin_adapters import (
    UnknownBinaryDocumentAdapter,
)
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.compatible_office_adapter import (
    CompatibleOfficeDocumentAdapter,
    NormalizedOfficeContainer,
    rebase_normalized_document_ir,
)
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.contract import DocumentSource
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.pipeline import (
    build_document_structure_ir,
)
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.planner import (
    plan_document_parsing,
)
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.registry import (
    DocumentAdapterRegistry,
)


class StaticNormalizer:
    name = "static-test-normalizer"
    version = "1"

    def __init__(self, outputs: dict[str, bytes]) -> None:
        self.outputs = dict(outputs)

    def available(self) -> bool:
        return True

    def normalize(self, source: DocumentSource, target_suffix: str) -> NormalizedOfficeContainer:
        data = self.outputs[target_suffix]
        filename = f"derived{target_suffix}"
        return NormalizedOfficeContainer(
            data=data,
            filename=filename,
            target_suffix=target_suffix,
            receipt={
                "schema": "qualibug.office-container-normalization-receipt.v1",
                "status": "COMPLETE",
                "normalizer_name": self.name,
                "normalizer_version": self.version,
                "source_filename": source.filename,
                "source_format": source.suffix.lstrip("."),
                "source_hash": source.content_hash,
                "target_format": target_suffix.lstrip("."),
                "derived_filename": filename,
                "derived_hash": hashlib.sha256(data).hexdigest(),
                "derived_container_is_not_evidence_root": True,
                "business_semantics_added": False,
            },
        )


def _xlsx_bytes() -> bytes:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "字段定义"
    sheet["A1"] = "字段"
    sheet["B1"] = "类型"
    sheet["A2"] = "order_id"
    sheet["B2"] = "bigint"
    buffer = io.BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def _docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("订单审批", level=1)
    document.add_paragraph("订单金额超过五万元需要财务审批")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "审批页面"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "规则"
    table.cell(1, 0).text = "金额"
    table.cell(1, 1).text = "大于50000显示财务审批"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.mark.parametrize(
    ("filename", "expected_target"),
    [
        ("需求.wps", ".docx"),
        ("模板.wpt", ".docx"),
        ("规则.et", ".xlsx"),
        ("字典.ett", ".xlsx"),
        ("旧表.xlsb", ".xlsx"),
        ("原型.dps", ".pptx"),
        ("演示模板.dpt", ".pptx"),
        ("旧需求.doc", ".docx"),
        ("旧表.xls", ".xlsx"),
        ("旧原型.ppt", ".pptx"),
    ],
)
def test_planner_selects_one_compatible_office_adapter(filename: str, expected_target: str) -> None:
    normalizer = StaticNormalizer(
        {".docx": b"docx", ".xlsx": b"xlsx", ".pptx": b"pptx"}
    )
    registry = DocumentAdapterRegistry(
        [CompatibleOfficeDocumentAdapter(normalizer), UnknownBinaryDocumentAdapter()]
    )
    source = DocumentSource("src", filename, b"legacy-container")

    plan = plan_document_parsing(source, registry)

    assert plan["status"] == "READY"
    assert plan["selected_adapters"][0]["adapter_name"] == "compatible-office-normalization"
    assert expected_target in plan["selected_adapters"][0]["reason"]
    assert plan["missing_capabilities"] == []


@pytest.mark.parametrize(
    ("filename", "target_suffix", "builder"),
    [
        ("需求.wps", ".docx", _docx_bytes),
        ("规则.et", ".xlsx", _xlsx_bytes),
        ("原型.dps", ".pptx", _pptx_bytes),
    ],
)
def test_compatible_office_pipeline_reuses_native_ir_and_original_evidence(
    filename: str,
    target_suffix: str,
    builder,
) -> None:
    original_bytes = ("original:" + filename).encode("utf-8")
    original_hash = hashlib.sha256(original_bytes).hexdigest()
    normalizer = StaticNormalizer({target_suffix: builder()})
    registry = DocumentAdapterRegistry(
        [CompatibleOfficeDocumentAdapter(normalizer), UnknownBinaryDocumentAdapter()]
    )

    result = build_document_structure_ir(
        original_bytes,
        filename=filename,
        source_id="src_enterprise_material",
        registry=registry,
    )

    assert result["format"] == filename.rsplit(".", 1)[1]
    assert result["office_normalization_receipt"]["derived_container_is_not_evidence_root"] is True
    assert result["office_normalization_receipt"]["source_hash"] == original_hash
    assert result["evidence_closure_receipt"]["source_hash"] == original_hash
    assert result["evidence_closure_receipt"]["source_traceability_rate"] == 1.0
    assert result["blocks"]
    assert all(block["source_hash"] == original_hash for block in result["blocks"])
    assert all(str(block["source_locator"]).startswith(filename) for block in result["blocks"])
    assert all(
        block["structure_evidence"]["container_normalization"][
            "derived_container_is_not_evidence_root"
        ]
        is True
        for block in result["blocks"]
    )
    reasons = {
        row.get("reason_code") for row in result.get("unsupported_content") or []
    }
    assert "OFFICE_COMPATIBILITY_CONTAINER_NORMALIZED" in reasons

    if target_suffix == ".xlsx":
        order_id = next(block for block in result["blocks"] if block.get("text") == "order_id")
        assert order_id["cell_ref"] == "A2"
        assert order_id["evidence_address"]["address_kind"] == "SPREADSHEET_CELL"
    if target_suffix == ".pptx":
        table_ids = {table["block_id"] for table in result["tables"]}
        table_cells = [block for block in result["blocks"] if block.get("type") == "TABLE_CELL"]
        assert table_cells
        assert all(block["parent_id"] in table_ids for block in table_cells)


def test_rebased_block_identity_does_not_depend_on_derived_container_hash() -> None:
    original = DocumentSource("src_rules", "规则.et", b"immutable-original")
    document_ir = {
        "schema": "qualibug.document-ir.v1",
        "format": "xlsx",
        "filename": "derived.xlsx",
        "plain_text": "状态\t规则",
        "blocks": [
            {
                "block_id": "derived-block",
                "type": "TABLE_CELL",
                "parent_id": "derived-table",
                "order": 1,
                "region": "body",
                "text": "状态",
                "source_locator": "derived.xlsx#sheet=规则;cell=A1",
                "sheet": "规则",
                "cell_ref": "A1",
                "structure_evidence": {"method": "openpyxl_native_cell"},
            }
        ],
        "sections": [],
        "tables": [
            {
                "block_id": "derived-table",
                "source_locator": "derived.xlsx#sheet=规则",
                "cell_block_ids": ["derived-block"],
            }
        ],
        "pages": [],
        "unsupported_content": [],
        "structure_receipt": {"status": "COMPLETE"},
    }

    first = rebase_normalized_document_ir(
        document_ir,
        original_source=original,
        normalized=NormalizedOfficeContainer(
            b"derived-one",
            "derived.xlsx",
            ".xlsx",
            {
                "normalizer_name": "test",
                "source_format": "et",
                "target_format": "xlsx",
                "derived_hash": "1" * 64,
            },
        ),
    )
    second = rebase_normalized_document_ir(
        document_ir,
        original_source=original,
        normalized=NormalizedOfficeContainer(
            b"derived-two",
            "derived.xlsx",
            ".xlsx",
            {
                "normalizer_name": "test",
                "source_format": "et",
                "target_format": "xlsx",
                "derived_hash": "2" * 64,
            },
        ),
    )

    assert first["blocks"][0]["block_id"] == second["blocks"][0]["block_id"]
    assert first["tables"][0]["block_id"] == second["tables"][0]["block_id"]
    assert first["blocks"][0]["source_locator"] == "规则.et#sheet=规则;cell=A1"
    assert first["blocks"][0]["parent_id"] == first["tables"][0]["block_id"]
