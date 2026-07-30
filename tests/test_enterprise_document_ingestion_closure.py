from __future__ import annotations

import io

from openpyxl import Workbook
from openpyxl.comments import Comment
from pptx import Presentation
from pptx.util import Inches

from ai_test_asset_center.enterprise_knowledge_center.document_ingestion import (
    build_document_structure_ir,
)
from ai_test_asset_center.enterprise_knowledge_center.enterprise_understanding.document_structure_gate import (
    apply_document_structure_completeness,
)
from ai_test_asset_center.enterprise_knowledge_center.enterprise_understanding.schema import (
    empty_model,
)


def _reason_codes(model: dict) -> set[str]:
    return {
        str(row.get("reason_code") or row.get("kind") or "")
        for row in model.get("unknowns") or []
        if isinstance(row, dict)
    }


def test_xlsx_uses_native_structure_and_exact_cell_evidence() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "工单规则"
    sheet["A1"] = "角色"
    sheet["B1"] = "规则"
    sheet["A2"] = "客服"
    sheet["B2"] = "客服必须关闭本人创建的工单"
    sheet["C2"] = "=1+1"
    sheet["B2"].comment = Comment("管理员除外", "产品经理")
    stream = io.BytesIO()
    workbook.save(stream)
    workbook.close()

    document_ir = build_document_structure_ir(
        stream.getvalue(),
        filename="ticket_rules.xlsx",
        source_id="source:xlsx",
    )

    assert document_ir["format"] == "xlsx"
    assert "spreadsheet-native-structure" in document_ir["adapter_merge_receipt"][
        "adapter_names"
    ]
    assert document_ir["evidence_closure_receipt"]["status"] == "PASS"
    assert document_ir["evidence_closure_receipt"]["source_traceability_rate"] == 1.0
    assert document_ir["evidence_closure_receipt"]["exact_address_rate"] == 1.0

    cells = [
        row
        for row in document_ir["blocks"]
        if row.get("type") == "TABLE_CELL"
    ]
    assert any(row.get("cell_ref") == "B2" for row in cells)
    assert any("#sheet=工单规则;cell=B2" in row.get("source_locator", "") for row in cells)
    assert all(row.get("source_hash") for row in cells)
    assert any(row.get("formula") == "=1+1" for row in cells)
    assert any(row.get("type") == "NOTE" and row.get("text") == "管理员除外" for row in document_ir["blocks"])


def test_pptx_uses_native_slide_shape_and_table_evidence() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    textbox.text_frame.text = "审批通过后，申请单进入已通过状态"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(7), Inches(2))
    table_shape.table.cell(0, 0).text = "角色"
    table_shape.table.cell(0, 1).text = "权限"
    table_shape.table.cell(1, 0).text = "审批人"
    table_shape.table.cell(1, 1).text = "可以审批"
    stream = io.BytesIO()
    presentation.save(stream)

    document_ir = build_document_structure_ir(
        stream.getvalue(),
        filename="approval_flow.pptx",
        source_id="source:pptx",
    )

    assert document_ir["format"] == "pptx"
    assert "presentation-native-structure" in document_ir["adapter_merge_receipt"][
        "adapter_names"
    ]
    assert document_ir["evidence_closure_receipt"]["status"] == "PASS"
    assert any(
        row.get("type") == "PARAGRAPH"
        and "#slide=1;shape=" in row.get("source_locator", "")
        for row in document_ir["blocks"]
    )
    assert any(
        row.get("type") == "TABLE_CELL"
        and ";table-cell=R2C2" in row.get("source_locator", "")
        for row in document_ir["blocks"]
    )
    assert all(
        row.get("source_hash")
        for row in document_ir["blocks"]
        if row.get("text")
    )


def test_text_ingestion_binds_every_formal_block_to_source_fingerprint() -> None:
    document_ir = build_document_structure_ir(
        "# 工单\n客服必须关闭本人创建的工单。".encode("utf-8"),
        filename="rules.md",
        source_id="source:text",
    )
    receipt = document_ir["evidence_closure_receipt"]
    assert receipt["status"] == "PASS"
    assert receipt["source_traceability_rate"] == 1.0
    assert receipt["untraceable_authority_block_count"] == 0
    assert all(
        row.get("source_id") == "source:text" and row.get("source_hash")
        for row in document_ir["blocks"]
        if row.get("text")
    )


def test_active_source_without_document_ir_blocks_understanding() -> None:
    model = empty_model()
    asset = {
        "source_inventory": [
            {"source_id": "source:missing", "status": "active"},
        ],
        "document_structure_assets": {
            "source_count": 0,
            "items": [],
            "errors": [],
        },
        "business_fact_ledger": {"items": []},
        "enterprise_comprehension_gate": {"entry_allowed": True},
    }

    result = apply_document_structure_completeness(model, asset)
    assert "ACTIVE_SOURCE_WITHOUT_DOCUMENT_STRUCTURE" in _reason_codes(result)
    assert result["gate"]["entry_allowed"] is False


def test_accepted_fact_without_exact_document_block_blocks_understanding() -> None:
    model = empty_model()
    structure = build_document_structure_ir(
        "# 工单\n客服可以查看工单。".encode("utf-8"),
        filename="rules.md",
        source_id="source:rules",
    )
    asset = {
        "source_inventory": [
            {"source_id": "source:rules", "status": "active"},
        ],
        "document_structure_assets": {
            "source_count": 1,
            "items": [{"source_id": "source:rules", **structure}],
            "errors": [],
        },
        "business_fact_ledger": {
            "items": [
                {
                    "fact_id": "fact:orphan",
                    "status": "ACCEPTED",
                    "kind": "RULE",
                }
            ]
        },
        "document_ir_fact_evidence_receipt": {
            "aligned_fact_count": 0,
            "unresolved_fact_count": 1,
            "aligned": [],
            "unresolved": [
                {
                    "fact_id": "fact:orphan",
                    "source_id": "source:rules",
                    "reason": "DOCUMENT_IR_FACT_BLOCK_NOT_FOUND",
                    "candidate_block_ids": [],
                }
            ],
        },
        "enterprise_comprehension_gate": {"entry_allowed": True},
    }

    result = apply_document_structure_completeness(model, asset)
    assert "FORMAL_FACT_WITHOUT_EXACT_DOCUMENT_EVIDENCE" in _reason_codes(result)
    assert result["document_structure_summary"][
        "accepted_fact_exact_evidence_rate"
    ] == 0.0
    assert result["gate"]["entry_allowed"] is False
