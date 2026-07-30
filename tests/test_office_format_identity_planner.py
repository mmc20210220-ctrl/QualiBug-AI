from __future__ import annotations

import io
import zipfile

import pytest

from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.contract import (
    DocumentSource,
)
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.pipeline import (
    build_document_structure_ir,
)
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.planner import (
    plan_document_parsing,
)
from ai_test_asset_center.enterprise_knowledge_center.document_ingestion.registry import (
    build_default_registry,
)


def _zip_members(*members: str) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types />")
        for member in members:
            archive.writestr(member, b"fixture")
    return buffer.getvalue()


@pytest.mark.parametrize(
    ("filename", "members", "expected_format", "expected_capability_family"),
    [
        ("需求.docm", ("word/document.xml",), "docm", "word"),
        ("模板.dotx", ("word/document.xml",), "dotx", "word"),
        ("规则.xlsm", ("xl/workbook.xml",), "xlsm", "spreadsheet"),
        ("二进制工作簿.xlsb", ("xl/workbook.bin",), "xlsb", "spreadsheet"),
        ("原型.pptm", ("ppt/presentation.xml",), "pptm", "presentation"),
        ("需求.wps", ("wps/content.xml",), "wps", "word"),
        ("规则.et", ("et/workbook.xml",), "et", "spreadsheet"),
        ("原型.dps", ("dps/presentation.xml",), "dps", "presentation"),
    ],
)
def test_planner_preserves_exact_office_source_format(
    filename: str,
    members: tuple[str, ...],
    expected_format: str,
    expected_capability_family: str,
) -> None:
    source = DocumentSource("src", filename, _zip_members(*members))

    plan = plan_document_parsing(source, build_default_registry())

    assert plan["detected_family"] == expected_format
    assert plan["detected_format"] == expected_format
    assert plan["capability_family"] == expected_capability_family
    assert plan["required_capabilities"]
    assert plan["selected_adapters"]
    assert plan["selected_adapters"][0]["adapter_name"] != "unknown-binary-fallback"


def _docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("订单审批", level=1)
    document.add_paragraph("超过五万元需要财务审批")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes() -> bytes:
    import openpyxl

    workbook = openpyxl.Workbook()
    workbook.active["A1"] = "订单编号"
    buffer = io.BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "审批页面"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.mark.parametrize(
    ("filename", "builder", "macro_reason"),
    [
        ("需求.docm", _docx_bytes, "WORD_MACRO_CODE_NOT_PARSED"),
        ("规则.xlsm", _xlsx_bytes, "SPREADSHEET_MACRO_CODE_NOT_PARSED"),
        ("原型.pptm", _pptx_bytes, "PRESENTATION_MACRO_CODE_NOT_PARSED"),
    ],
)
def test_pipeline_keeps_macro_subtype_after_native_structure_parse(
    filename: str,
    builder,
    macro_reason: str,
) -> None:
    result = build_document_structure_ir(
        builder(),
        filename=filename,
        source_id="src_macro_identity",
    )

    assert result["format"] == filename.rsplit(".", 1)[1]
    assert result["structure_receipt"]["format"] == result["format"]
    assert result["parsing_plan"]["detected_format"] == result["format"]
    assert result["blocks"]
    reasons = {
        row.get("reason_code") for row in result.get("unsupported_content") or []
    }
    assert macro_reason in reasons
