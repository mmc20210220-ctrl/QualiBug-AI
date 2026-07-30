"""Native OOXML spreadsheet and presentation adapters.

These adapters preserve container structure and exact source addresses.  They do not
infer business meaning from sheet order, slide order, styling, formulas, or diagrams.
Unsupported embedded content remains fail-visible in ``unsupported_content``.
"""
from __future__ import annotations

import hashlib
import io
import posixpath
import zipfile
from collections import Counter
from typing import Any, Iterable
from xml.etree import ElementTree

from .._document_structure_ir import DOCUMENT_IR_SCHEMA, STRUCTURE_RECEIPT_SCHEMA
from .contract import (
    AdapterMatch,
    CAP_COMMENT_EXTRACTION,
    CAP_FORMULA_EXTRACTION,
    CAP_HEADING_HIERARCHY,
    CAP_IMAGE_PRESENCE,
    CAP_STYLE_SEMANTICS,
    CAP_TABLE_STRUCTURE,
    CAP_TEXT_EXTRACTION,
    DocumentAdapter,
    DocumentSource,
    MODE_PRIMARY,
)

_MAX_SPREADSHEET_CELLS = 500_000
_MAX_PRESENTATION_SHAPES = 50_000
_MAX_CELL_TEXT = 20_000


def _text(value: Any) -> str:
    return str(value if value is not None else "").strip()


def _stable_id(prefix: str, *parts: Any) -> str:
    material = "\x1f".join(_text(part) for part in parts)
    return f"{prefix}:{hashlib.sha256(material.encode('utf-8')).hexdigest()[:20]}"


def _safe_locator_token(value: Any) -> str:
    return _text(value).replace("#", "%23").replace(";", "%3B")


def _looks_like_ooxml(data: bytes, member: str) -> bool:
    if not data.startswith(b"PK"):
        return False
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = set(archive.namelist())
            return member in names and "[Content_Types].xml" in names
    except Exception:
        return False


def _dedupe_gaps(rows: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    result: dict[tuple[str, str], dict[str, Any]] = {}
    for raw in rows:
        if not isinstance(raw, dict):
            continue
        row = dict(raw)
        kind = _text(row.get("kind") or row.get("reason_code") or "OFFICE_CONTENT_UNPARSED")
        locator = _text(row.get("source_locator"))
        key = (kind, locator)
        if key not in result:
            result[key] = row
            continue
        result[key]["count"] = int(result[key].get("count") or 0) + int(row.get("count") or 0)
        result[key]["blocks_formal_understanding"] = bool(
            result[key].get("blocks_formal_understanding")
        ) or bool(row.get("blocks_formal_understanding"))
    return list(result.values())


def _structure_receipt(
    *,
    source_format: str,
    blocks: list[dict[str, Any]],
    sections: list[dict[str, Any]],
    tables: list[dict[str, Any]],
    unsupported: list[dict[str, Any]],
    extra: dict[str, Any] | None = None,
) -> dict[str, Any]:
    critical = sum(
        int(row.get("count") or 0)
        for row in unsupported
        if bool(row.get("blocks_formal_understanding"))
    )
    status = "BLOCKED" if critical else "PARTIAL" if unsupported else "COMPLETE"
    receipt = {
        "schema": STRUCTURE_RECEIPT_SCHEMA,
        "status": status,
        "format": source_format,
        "block_count": len(blocks),
        "source_traceability_rate": round(
            sum(1 for row in blocks if _text(row.get("source_locator"))) / len(blocks), 4
        )
        if blocks
        else 1.0,
        "block_type_distribution": dict(Counter(_text(row.get("type")) for row in blocks)),
        "section_count": len(sections),
        "table_count": len(tables),
        "unsupported_content_count": sum(int(row.get("count") or 0) for row in unsupported),
        "critical_unsupported_content_count": critical,
        "unsupported_content": unsupported,
        "document_order_is_business_flow": False,
        "filename_is_business_context": False,
    }
    receipt.update(dict(extra or {}))
    return receipt


def _cell_style(cell: Any) -> dict[str, Any]:
    font = getattr(cell, "font", None)
    fill = getattr(cell, "fill", None)
    alignment = getattr(cell, "alignment", None)
    border = getattr(cell, "border", None)
    return {
        "style_id": getattr(cell, "style_id", None),
        "number_format": _text(getattr(cell, "number_format", "")),
        "font": {
            "name": _text(getattr(font, "name", "")),
            "size": getattr(getattr(font, "sz", None), "real", getattr(font, "sz", None)),
            "bold": bool(getattr(font, "bold", False)),
            "italic": bool(getattr(font, "italic", False)),
        },
        "fill_type": _text(getattr(fill, "fill_type", "")),
        "horizontal_alignment": _text(getattr(alignment, "horizontal", "")),
        "vertical_alignment": _text(getattr(alignment, "vertical", "")),
        "border_present": any(
            _text(getattr(getattr(border, edge, None), "style", ""))
            for edge in ("left", "right", "top", "bottom")
        ),
        "semantic_interpretation_applied": False,
    }


def _defined_names(workbook: Any) -> list[Any]:
    container = getattr(workbook, "defined_names", None)
    if container is None:
        return []
    values = getattr(container, "values", None)
    if callable(values):
        try:
            return list(values())
        except Exception:
            pass
    legacy = getattr(container, "definedName", None)
    return list(legacy or [])


class SpreadsheetDocumentAdapter(DocumentAdapter):
    """Preserve workbook/sheet/cell/formula/comment structure for OOXML spreadsheets."""

    name = "spreadsheet-native-structure"
    parser_version = "1"
    priority = 110
    mode = MODE_PRIMARY
    capabilities = frozenset(
        {
            CAP_TEXT_EXTRACTION,
            CAP_TABLE_STRUCTURE,
            CAP_FORMULA_EXTRACTION,
            CAP_COMMENT_EXTRACTION,
            CAP_STYLE_SEMANTICS,
        }
    )
    _SUFFIXES = {".xlsx", ".xlsm", ".xltx", ".xltm"}

    def probe(self, source: DocumentSource) -> AdapterMatch | None:
        signature = _looks_like_ooxml(source.data, "xl/workbook.xml")
        suffix = source.suffix in self._SUFFIXES
        if not signature and not suffix:
            return None
        return AdapterMatch(
            self.name,
            125 if signature else 96,
            "office_open_xml_excel_container" if signature else "spreadsheet_filename_suffix",
            tuple(sorted(self.capabilities)),
            self.mode,
        )

    def extract(self, source: DocumentSource) -> dict[str, Any]:
        import openpyxl

        keep_vba = source.suffix in {".xlsm", ".xltm"}
        workbook = openpyxl.load_workbook(
            io.BytesIO(source.data),
            read_only=False,
            data_only=False,
            keep_vba=keep_vba,
            keep_links=True,
        )
        blocks: list[dict[str, Any]] = []
        sections: list[dict[str, Any]] = []
        tables: list[dict[str, Any]] = []
        unsupported: list[dict[str, Any]] = []
        plain_lines: list[str] = []
        order = 0
        processed_cells = 0
        formula_count = 0
        comment_count = 0
        validation_count = 0
        defined_name_count = 0
        merged_range_count = 0
        image_count = 0
        chart_count = 0

        try:
            for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
                sheet_name = _text(worksheet.title) or f"Sheet{sheet_index}"
                sheet_token = _safe_locator_token(sheet_name)
                sheet_locator = f"{source.filename}#sheet={sheet_token}"
                order += 1
                heading_id = _stable_id(
                    "spreadsheet_sheet", source.source_id, source.content_hash, sheet_index, sheet_name
                )
                heading = {
                    "block_id": heading_id,
                    "type": "HEADING",
                    "parent_id": "",
                    "order": order,
                    "region": "body",
                    "level": 1,
                    "text": sheet_name,
                    "source_locator": sheet_locator,
                    "sheet": sheet_name,
                    "sheet_index": sheet_index,
                    "sheet_state": _text(getattr(worksheet, "sheet_state", "visible")),
                    "structure_evidence": {"method": "native_workbook_sheet_identity"},
                }
                blocks.append(heading)
                sections.append(
                    {
                        "block_id": heading_id,
                        "level": 1,
                        "title": sheet_name,
                        "source_locator": sheet_locator,
                        "sheet": sheet_name,
                        "sheet_index": sheet_index,
                    }
                )
                plain_lines.append(f"## {sheet_name}")

                merged_lookup: dict[str, str] = {}
                for merged in worksheet.merged_cells.ranges:
                    merged_range_count += 1
                    for row in worksheet[merged.coord]:
                        for cell in row:
                            merged_lookup[cell.coordinate] = merged.coord

                sheet_cell_ids: list[str] = []
                first_nonempty_row: int | None = None
                row_text: dict[int, list[tuple[int, str]]] = {}
                stop = False
                for row in worksheet.iter_rows():
                    for cell in row:
                        if cell.value is None and getattr(cell, "comment", None) is None:
                            continue
                        processed_cells += 1
                        if processed_cells > _MAX_SPREADSHEET_CELLS:
                            unsupported.append(
                                {
                                    "kind": "SPREADSHEET_CELL_LIMIT_EXCEEDED",
                                    "reason_code": "SPREADSHEET_CELL_LIMIT_EXCEEDED",
                                    "count": 1,
                                    "status": "SOURCE_NOT_FULLY_PARSED",
                                    "severity": "P0",
                                    "blocks_formal_understanding": True,
                                    "included_in_plain_text_authority": False,
                                    "source_locator": sheet_locator,
                                    "cell_limit": _MAX_SPREADSHEET_CELLS,
                                }
                            )
                            stop = True
                            break
                        coordinate = _text(cell.coordinate)
                        value = _text(cell.value)
                        if len(value) > _MAX_CELL_TEXT:
                            unsupported.append(
                                {
                                    "kind": "SPREADSHEET_CELL_TEXT_TRUNCATED",
                                    "reason_code": "SPREADSHEET_CELL_TEXT_TRUNCATED",
                                    "count": 1,
                                    "status": "CELL_TEXT_EXCEEDS_SAFE_LIMIT",
                                    "severity": "P0",
                                    "blocks_formal_understanding": True,
                                    "included_in_plain_text_authority": False,
                                    "source_locator": f"{sheet_locator};cell={coordinate}",
                                    "original_length": len(value),
                                    "safe_limit": _MAX_CELL_TEXT,
                                }
                            )
                            value = value[:_MAX_CELL_TEXT]
                        formula = value if _text(getattr(cell, "data_type", "")) == "f" else ""
                        if formula:
                            formula_count += 1
                        locator = f"{sheet_locator};cell={coordinate}"
                        order += 1
                        block_id = _stable_id(
                            "spreadsheet_cell",
                            source.source_id,
                            source.content_hash,
                            sheet_name,
                            coordinate,
                        )
                        block = {
                            "block_id": block_id,
                            "type": "TABLE_CELL",
                            "parent_id": heading_id,
                            "order": order,
                            "region": "body",
                            "text": value,
                            "source_locator": locator,
                            "sheet": sheet_name,
                            "sheet_index": sheet_index,
                            "cell_ref": coordinate,
                            "row_index": int(cell.row),
                            "column_index": int(cell.column),
                            "formula": formula,
                            "data_type": _text(getattr(cell, "data_type", "")),
                            "merged_range": merged_lookup.get(coordinate, ""),
                            "hyperlink": _text(getattr(getattr(cell, "hyperlink", None), "target", "")),
                            "style_evidence": _cell_style(cell),
                            "structure_evidence": {
                                "method": "openpyxl_native_cell",
                                "source_backed": True,
                            },
                        }
                        blocks.append(block)
                        sheet_cell_ids.append(block_id)
                        if first_nonempty_row is None:
                            first_nonempty_row = int(cell.row)
                        if value:
                            row_text.setdefault(int(cell.row), []).append((int(cell.column), value))

                        comment = getattr(cell, "comment", None)
                        comment_text = _text(getattr(comment, "text", ""))
                        if comment_text:
                            comment_count += 1
                            order += 1
                            comment_locator = f"{locator};comment"
                            blocks.append(
                                {
                                    "block_id": _stable_id(
                                        "spreadsheet_comment",
                                        source.source_id,
                                        source.content_hash,
                                        sheet_name,
                                        coordinate,
                                        comment_text,
                                    ),
                                    "type": "NOTE",
                                    "parent_id": block_id,
                                    "order": order,
                                    "region": "body",
                                    "text": comment_text,
                                    "source_locator": comment_locator,
                                    "sheet": sheet_name,
                                    "cell_ref": coordinate,
                                    "author": _text(getattr(comment, "author", "")),
                                    "structure_evidence": {
                                        "method": "openpyxl_native_comment",
                                        "source_backed": True,
                                    },
                                }
                            )
                    if stop:
                        break
                for row_index in sorted(row_text):
                    values = [value for _column, value in sorted(row_text[row_index])]
                    if values:
                        plain_lines.append("\t".join(values))
                plain_lines.append("")

                table_id = _stable_id(
                    "spreadsheet_table", source.source_id, source.content_hash, sheet_name
                )
                tables.append(
                    {
                        "block_id": table_id,
                        "type": "TABLE",
                        "source_locator": sheet_locator,
                        "sheet": sheet_name,
                        "sheet_index": sheet_index,
                        "cell_block_ids": sheet_cell_ids,
                        "first_nonempty_row_candidate": first_nonempty_row,
                        "header_semantics_confirmed": False,
                        "document_order_is_business_flow": False,
                    }
                )

                validations = list(
                    getattr(getattr(worksheet, "data_validations", None), "dataValidation", [])
                    or []
                )
                for validation_index, validation in enumerate(validations, start=1):
                    validation_count += 1
                    validation_text = "; ".join(
                        part
                        for part in [
                            f"range={_text(getattr(validation, 'sqref', ''))}",
                            f"type={_text(getattr(validation, 'type', ''))}",
                            f"operator={_text(getattr(validation, 'operator', ''))}",
                            f"formula1={_text(getattr(validation, 'formula1', ''))}",
                            f"formula2={_text(getattr(validation, 'formula2', ''))}",
                            f"allow_blank={bool(getattr(validation, 'allow_blank', False))}",
                        ]
                        if part.split("=", 1)[1] not in {"", "False"}
                    )
                    if not validation_text:
                        continue
                    order += 1
                    locator = f"{sheet_locator};data-validation={validation_index}"
                    blocks.append(
                        {
                            "block_id": _stable_id(
                                "spreadsheet_data_validation",
                                source.source_id,
                                source.content_hash,
                                sheet_name,
                                validation_index,
                                validation_text,
                            ),
                            "type": "FORMULA",
                            "parent_id": heading_id,
                            "order": order,
                            "region": "body",
                            "text": validation_text,
                            "source_locator": locator,
                            "sheet": sheet_name,
                            "validation_index": validation_index,
                            "structure_evidence": {
                                "method": "openpyxl_native_data_validation",
                                "source_backed": True,
                                "business_semantics_inferred": False,
                            },
                        }
                    )
                    plain_lines.append(validation_text)

                sheet_images = list(getattr(worksheet, "_images", []) or [])
                sheet_charts = list(getattr(worksheet, "_charts", []) or [])
                image_count += len(sheet_images)
                chart_count += len(sheet_charts)
                if sheet_images:
                    unsupported.append(
                        {
                            "kind": "SPREADSHEET_EMBEDDED_IMAGE_NOT_SEMANTICALLY_PARSED",
                            "reason_code": "SPREADSHEET_EMBEDDED_IMAGE_NOT_SEMANTICALLY_PARSED",
                            "count": len(sheet_images),
                            "status": "IMAGE_PRESENCE_RECORDED_CONTENT_UNPARSED",
                            "severity": "P1",
                            "blocks_formal_understanding": False,
                            "included_in_plain_text_authority": False,
                            "source_locator": sheet_locator,
                        }
                    )
                if sheet_charts:
                    unsupported.append(
                        {
                            "kind": "SPREADSHEET_CHART_STRUCTURE_NOT_PARSED",
                            "reason_code": "SPREADSHEET_CHART_STRUCTURE_NOT_PARSED",
                            "count": len(sheet_charts),
                            "status": "CHART_PRESENCE_RECORDED_CONTENT_UNPARSED",
                            "severity": "P1",
                            "blocks_formal_understanding": False,
                            "included_in_plain_text_authority": False,
                            "source_locator": sheet_locator,
                        }
                    )
                if stop:
                    break

            for index, defined_name in enumerate(_defined_names(workbook), start=1):
                name = _text(getattr(defined_name, "name", ""))
                expression = _text(getattr(defined_name, "attr_text", ""))
                if not name and not expression:
                    continue
                defined_name_count += 1
                order += 1
                locator = f"{source.filename}#defined-name={_safe_locator_token(name or index)}"
                value = f"{name}={expression}" if name else expression
                blocks.append(
                    {
                        "block_id": _stable_id(
                            "spreadsheet_defined_name",
                            source.source_id,
                            source.content_hash,
                            name,
                            expression,
                        ),
                        "type": "FORMULA",
                        "parent_id": "",
                        "order": order,
                        "region": "body",
                        "text": value,
                        "source_locator": locator,
                        "defined_name": name,
                        "formula": expression,
                        "structure_evidence": {
                            "method": "openpyxl_native_defined_name",
                            "source_backed": True,
                        },
                    }
                )
                plain_lines.append(value)

            if keep_vba:
                unsupported.append(
                    {
                        "kind": "SPREADSHEET_MACRO_CODE_NOT_PARSED",
                        "reason_code": "SPREADSHEET_MACRO_CODE_NOT_PARSED",
                        "count": 1,
                        "status": "MACRO_CONTAINER_PRESENT_CODE_UNAVAILABLE",
                        "severity": "P0",
                        "blocks_formal_understanding": True,
                        "included_in_plain_text_authority": False,
                        "source_locator": f"{source.filename}#vba-project",
                    }
                )
        finally:
            workbook.close()

        unsupported = _dedupe_gaps(unsupported)
        receipt = _structure_receipt(
            source_format=source.suffix.lstrip(".") or "xlsx",
            blocks=blocks,
            sections=sections,
            tables=tables,
            unsupported=unsupported,
            extra={
                "sheet_count": len(sections),
                "spreadsheet_cell_count": processed_cells,
                "formula_cell_count": formula_count,
                "comment_count": comment_count,
                "data_validation_count": validation_count,
                "defined_name_count": defined_name_count,
                "merged_range_count": merged_range_count,
                "image_count": image_count,
                "chart_count": chart_count,
                "native_spreadsheet_structure": True,
                "style_metadata_is_evidence_not_business_semantics": True,
            },
        )
        return {
            "schema": DOCUMENT_IR_SCHEMA,
            "format": source.suffix.lstrip(".") or "xlsx",
            "filename": source.filename,
            "plain_text": "\n".join(plain_lines).strip(),
            "blocks": blocks,
            "sections": sections,
            "tables": tables,
            "pages": [],
            "unsupported_content": unsupported,
            "structure_receipt": receipt,
        }


def _presentation_notes_by_slide(data: bytes) -> dict[int, str]:
    """Read speaker notes from OOXML relationships without mutating the package."""
    result: dict[int, str] = {}
    relationship_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    drawing_text = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = set(archive.namelist())
            slide_index = 1
            while f"ppt/slides/slide{slide_index}.xml" in names:
                rel_path = f"ppt/slides/_rels/slide{slide_index}.xml.rels"
                if rel_path not in names:
                    slide_index += 1
                    continue
                rel_root = ElementTree.fromstring(archive.read(rel_path))
                target = ""
                for relationship in rel_root.findall("r:Relationship", relationship_ns):
                    if _text(relationship.attrib.get("Type")).endswith("/notesSlide"):
                        target = _text(relationship.attrib.get("Target"))
                        break
                if target:
                    resolved = posixpath.normpath(
                        posixpath.join(posixpath.dirname(f"ppt/slides/slide{slide_index}.xml"), target)
                    )
                    if resolved in names:
                        notes_root = ElementTree.fromstring(archive.read(resolved))
                        values = [_text(node.text) for node in notes_root.iter(drawing_text) if _text(node.text)]
                        if values:
                            result[slide_index] = "\n".join(values)
                slide_index += 1
    except Exception:
        return {}
    return result


def _iter_shapes(shapes: Any, *, parent_shape_id: str = "") -> Iterable[tuple[Any, str]]:
    for shape in shapes:
        shape_id = _text(getattr(shape, "shape_id", ""))
        composite = f"{parent_shape_id}.{shape_id}" if parent_shape_id else shape_id
        yield shape, composite
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from _iter_shapes(nested, parent_shape_id=composite)


class PresentationDocumentAdapter(DocumentAdapter):
    """Preserve slide/shape/paragraph/table/note structure for OOXML presentations."""

    name = "presentation-native-structure"
    parser_version = "1"
    priority = 110
    mode = MODE_PRIMARY
    capabilities = frozenset(
        {
            CAP_TEXT_EXTRACTION,
            CAP_HEADING_HIERARCHY,
            CAP_TABLE_STRUCTURE,
            CAP_IMAGE_PRESENCE,
            CAP_STYLE_SEMANTICS,
            CAP_COMMENT_EXTRACTION,
        }
    )
    _SUFFIXES = {".pptx", ".pptm", ".potx", ".potm", ".ppsx", ".ppsm"}

    def probe(self, source: DocumentSource) -> AdapterMatch | None:
        signature = _looks_like_ooxml(source.data, "ppt/presentation.xml")
        suffix = source.suffix in self._SUFFIXES
        if not signature and not suffix:
            return None
        return AdapterMatch(
            self.name,
            125 if signature else 96,
            "office_open_xml_presentation_container" if signature else "presentation_filename_suffix",
            tuple(sorted(self.capabilities)),
            self.mode,
        )

    def extract(self, source: DocumentSource) -> dict[str, Any]:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(source.data))
        notes_by_slide = _presentation_notes_by_slide(source.data)
        blocks: list[dict[str, Any]] = []
        sections: list[dict[str, Any]] = []
        tables: list[dict[str, Any]] = []
        pages: list[dict[str, Any]] = []
        unsupported: list[dict[str, Any]] = []
        plain_lines: list[str] = []
        order = 0
        shape_count = 0
        text_shape_count = 0
        table_count = 0
        table_cell_count = 0
        image_count = 0
        chart_count = 0
        diagram_count = 0
        notes_count = 0

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_locator = f"{source.filename}#slide={slide_index}"
            title = ""
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape is not None and getattr(title_shape, "has_text_frame", False):
                title = _text(getattr(title_shape, "text", ""))
            heading_text = title or f"Slide {slide_index}"
            order += 1
            heading_id = _stable_id(
                "presentation_slide", source.source_id, source.content_hash, slide_index
            )
            blocks.append(
                {
                    "block_id": heading_id,
                    "type": "HEADING",
                    "parent_id": "",
                    "order": order,
                    "region": "body",
                    "level": 1,
                    "text": heading_text,
                    "source_locator": slide_locator,
                    "slide": slide_index,
                    "structure_evidence": {"method": "native_presentation_slide_identity"},
                }
            )
            sections.append(
                {
                    "block_id": heading_id,
                    "level": 1,
                    "title": heading_text,
                    "source_locator": slide_locator,
                    "slide": slide_index,
                }
            )
            pages.append(
                {
                    "page": slide_index,
                    "slide": slide_index,
                    "source_locator": slide_locator,
                    "title": title,
                }
            )
            plain_lines.append(f"## {heading_text}")
            slide_text_found = bool(title)
            slide_image_count = 0

            for shape, composite_shape_id in _iter_shapes(slide.shapes):
                shape_count += 1
                if shape_count > _MAX_PRESENTATION_SHAPES:
                    unsupported.append(
                        {
                            "kind": "PRESENTATION_SHAPE_LIMIT_EXCEEDED",
                            "reason_code": "PRESENTATION_SHAPE_LIMIT_EXCEEDED",
                            "count": 1,
                            "status": "SOURCE_NOT_FULLY_PARSED",
                            "severity": "P0",
                            "blocks_formal_understanding": True,
                            "included_in_plain_text_authority": False,
                            "source_locator": slide_locator,
                            "shape_limit": _MAX_PRESENTATION_SHAPES,
                        }
                    )
                    break
                shape_locator = f"{slide_locator};shape={_safe_locator_token(composite_shape_id)}"
                has_table = bool(getattr(shape, "has_table", False))
                has_chart = bool(getattr(shape, "has_chart", False))
                has_text_frame = bool(getattr(shape, "has_text_frame", False))
                shape_type = _text(getattr(getattr(shape, "shape_type", None), "name", getattr(shape, "shape_type", "")))

                if has_text_frame:
                    paragraph_rows: list[str] = []
                    text_frame = getattr(shape, "text_frame", None)
                    for paragraph_index, paragraph in enumerate(
                        list(getattr(text_frame, "paragraphs", []) or []), start=1
                    ):
                        value = _text(getattr(paragraph, "text", ""))
                        if not value:
                            continue
                        slide_text_found = True
                        text_shape_count += 1
                        paragraph_rows.append(value)
                        order += 1
                        paragraph_locator = f"{shape_locator};paragraph={paragraph_index}"
                        blocks.append(
                            {
                                "block_id": _stable_id(
                                    "presentation_paragraph",
                                    source.source_id,
                                    source.content_hash,
                                    slide_index,
                                    composite_shape_id,
                                    paragraph_index,
                                    value,
                                ),
                                "type": "PARAGRAPH",
                                "parent_id": heading_id,
                                "order": order,
                                "region": "body",
                                "text": value,
                                "source_locator": paragraph_locator,
                                "slide": slide_index,
                                "shape_id": composite_shape_id,
                                "paragraph_index": paragraph_index,
                                "shape_type": shape_type,
                                "level": int(getattr(paragraph, "level", 0) or 0),
                                "bbox": {
                                    "left": int(getattr(shape, "left", 0) or 0),
                                    "top": int(getattr(shape, "top", 0) or 0),
                                    "width": int(getattr(shape, "width", 0) or 0),
                                    "height": int(getattr(shape, "height", 0) or 0),
                                    "unit": "EMU",
                                },
                                "structure_evidence": {
                                    "method": "python_pptx_native_text_frame",
                                    "source_backed": True,
                                },
                            }
                        )
                    plain_lines.extend(paragraph_rows)

                if has_table:
                    table_count += 1
                    table = shape.table
                    table_id = _stable_id(
                        "presentation_table",
                        source.source_id,
                        source.content_hash,
                        slide_index,
                        composite_shape_id,
                    )
                    cell_ids: list[str] = []
                    for row_index, row in enumerate(table.rows, start=1):
                        row_values: list[str] = []
                        for column_index, cell in enumerate(row.cells, start=1):
                            value = _text(cell.text)
                            if not value:
                                continue
                            slide_text_found = True
                            table_cell_count += 1
                            row_values.append(value)
                            order += 1
                            locator = (
                                f"{shape_locator};table-cell=R{row_index}C{column_index}"
                            )
                            block_id = _stable_id(
                                "presentation_table_cell",
                                source.source_id,
                                source.content_hash,
                                slide_index,
                                composite_shape_id,
                                row_index,
                                column_index,
                            )
                            cell_ids.append(block_id)
                            blocks.append(
                                {
                                    "block_id": block_id,
                                    "type": "TABLE_CELL",
                                    "parent_id": table_id,
                                    "order": order,
                                    "region": "body",
                                    "text": value,
                                    "source_locator": locator,
                                    "slide": slide_index,
                                    "shape_id": composite_shape_id,
                                    "row_index": row_index,
                                    "column_index": column_index,
                                    "cell_ref": f"R{row_index}C{column_index}",
                                    "bbox": {
                                        "left": int(getattr(shape, "left", 0) or 0),
                                        "top": int(getattr(shape, "top", 0) or 0),
                                        "width": int(getattr(shape, "width", 0) or 0),
                                        "height": int(getattr(shape, "height", 0) or 0),
                                        "unit": "EMU",
                                    },
                                    "structure_evidence": {
                                        "method": "python_pptx_native_table_cell",
                                        "source_backed": True,
                                    },
                                }
                            )
                        if row_values:
                            plain_lines.append("\t".join(row_values))
                    tables.append(
                        {
                            "block_id": table_id,
                            "type": "TABLE",
                            "source_locator": shape_locator,
                            "slide": slide_index,
                            "shape_id": composite_shape_id,
                            "cell_block_ids": cell_ids,
                            "header_semantics_confirmed": False,
                            "document_order_is_business_flow": False,
                        }
                    )

                if has_chart:
                    chart_count += 1
                    unsupported.append(
                        {
                            "kind": "PRESENTATION_CHART_STRUCTURE_NOT_PARSED",
                            "reason_code": "PRESENTATION_CHART_STRUCTURE_NOT_PARSED",
                            "count": 1,
                            "status": "CHART_PRESENCE_RECORDED_CONTENT_UNPARSED",
                            "severity": "P1",
                            "blocks_formal_understanding": False,
                            "included_in_plain_text_authority": False,
                            "source_locator": shape_locator,
                        }
                    )
                if shape_type in {"PICTURE", "LINKED_PICTURE"}:
                    image_count += 1
                    slide_image_count += 1
                if shape_type == "GRAPHIC_FRAME" and not has_table and not has_chart:
                    diagram_count += 1
                    unsupported.append(
                        {
                            "kind": "PRESENTATION_DIAGRAM_STRUCTURE_NOT_PARSED",
                            "reason_code": "PRESENTATION_DIAGRAM_STRUCTURE_NOT_PARSED",
                            "count": 1,
                            "status": "DIAGRAM_PRESENCE_RECORDED_CONTENT_UNPARSED",
                            "severity": "P1",
                            "blocks_formal_understanding": False,
                            "included_in_plain_text_authority": False,
                            "source_locator": shape_locator,
                        }
                    )

            notes = _text(notes_by_slide.get(slide_index))
            if notes:
                notes_count += 1
                order += 1
                blocks.append(
                    {
                        "block_id": _stable_id(
                            "presentation_notes",
                            source.source_id,
                            source.content_hash,
                            slide_index,
                            notes,
                        ),
                        "type": "NOTE",
                        "parent_id": heading_id,
                        "order": order,
                        "region": "body",
                        "text": notes,
                        "source_locator": f"{slide_locator};speaker-notes",
                        "slide": slide_index,
                        "structure_evidence": {
                            "method": "ooxml_notes_slide_text",
                            "source_backed": True,
                        },
                    }
                )
                plain_lines.append(notes)
                slide_text_found = True
            if slide_image_count:
                unsupported.append(
                    {
                        "kind": "PRESENTATION_IMAGE_CONTENT_UNPARSED",
                        "reason_code": "PRESENTATION_IMAGE_CONTENT_UNPARSED",
                        "count": slide_image_count,
                        "status": "IMAGE_PRESENCE_RECORDED_CONTENT_UNPARSED",
                        "severity": "P0" if not slide_text_found else "P1",
                        "blocks_formal_understanding": not slide_text_found,
                        "included_in_plain_text_authority": False,
                        "source_locator": slide_locator,
                        "image_only_slide": not slide_text_found,
                    }
                )
            plain_lines.append("")
            if shape_count > _MAX_PRESENTATION_SHAPES:
                break

        if source.suffix in {".pptm", ".potm", ".ppsm"}:
            unsupported.append(
                {
                    "kind": "PRESENTATION_MACRO_CODE_NOT_PARSED",
                    "reason_code": "PRESENTATION_MACRO_CODE_NOT_PARSED",
                    "count": 1,
                    "status": "MACRO_CONTAINER_PRESENT_CODE_UNAVAILABLE",
                    "severity": "P0",
                    "blocks_formal_understanding": True,
                    "included_in_plain_text_authority": False,
                    "source_locator": f"{source.filename}#vba-project",
                }
            )

        unsupported = _dedupe_gaps(unsupported)
        receipt = _structure_receipt(
            source_format=source.suffix.lstrip(".") or "pptx",
            blocks=blocks,
            sections=sections,
            tables=tables,
            unsupported=unsupported,
            extra={
                "page_count": len(pages),
                "slide_count": len(pages),
                "presentation_shape_count": shape_count,
                "text_shape_count": text_shape_count,
                "presentation_table_count": table_count,
                "presentation_table_cell_count": table_cell_count,
                "image_count": image_count,
                "chart_count": chart_count,
                "diagram_count": diagram_count,
                "speaker_notes_count": notes_count,
                "native_presentation_structure": True,
                "slide_order_is_not_business_flow": True,
            },
        )
        return {
            "schema": DOCUMENT_IR_SCHEMA,
            "format": source.suffix.lstrip(".") or "pptx",
            "filename": source.filename,
            "plain_text": "\n".join(plain_lines).strip(),
            "blocks": blocks,
            "sections": sections,
            "tables": tables,
            "pages": pages,
            "unsupported_content": unsupported,
            "structure_receipt": receipt,
        }


__all__ = ["SpreadsheetDocumentAdapter", "PresentationDocumentAdapter"]
